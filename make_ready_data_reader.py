from re import L
import pandas as pd
import matplotlib.pyplot as plt
import math
import os
import shutil

import datetime
import seaborn as sns
import matplotlib.pyplot as plt
# from matplotlib import dates
import matplotlib.dates as mdates

from datetime import timedelta
import numpy as np
import plotly.graph_objects as go

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
from pptx.util import Inches



# from datetime import datetime

## Change stuff here
# long_dir = r'C:\Users\dean.huang\main\projects\HeidelBergTimeTracking\MR Time Improvement Dean Edit'
long_dir = r'C:\Users\dean.huang\main\projects\HeidelBergTimeTracking\todd_build'
current = os.path.join(long_dir, "kpi_master_brett.xlsx")
# This will hold all plants data, each plant will be one dict
plant_dict = {}
root_die_historic_dict = {}
current_month_dict = {}
# Monthly - Plant dictionary
mr_month_dict = {}
month_label = ['Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
color_list = ['steelblue', 'firebrick', 'seagreen','mediumblue']


#Stuff to Change!
plant_list = ['Staunton', 'Fort Smith', 'Phoenixville', 'Oroville', 'Mississauga', 'Elk Grove Village', 'Sturgis', 'Vancouver']
# plant_list = ['Staunton', 'Fort Smith', 'Phoenixville']
# plant_list = ['Staunton']


save_fig = 1
print_graphs = 0



def multi_sheet_excel_writer():
    df1 = pd.DataFrame({'Data': ['a', 'b', 'c', 'd'], 'Data2': ['test1', 'test2','','']})
    df2 = pd.DataFrame({'Data': [1, 2, 3, 4]})
    df3 = pd.DataFrame({'Data': [1.1, 1.2, 1.3, 1.4]})
    df4 = pd.DataFrame({'Data': [1.1, 1.2, 1.3, 1.4]})
    writer = pd.ExcelWriter('multiple.xlsx', engine='xlsxwriter')
    df1.to_excel(writer, sheet_name='Sheeta')
    df2.to_excel(writer, sheet_name='Sheetb')
    df3.to_excel(writer, sheet_name='Sheetc')
    df4.to_excel(writer, sheet_name='Sheetd')
    writer.save()

def save_to_excel(filename, data_dict):
    combined_df = pd.DataFrame.from_dict(data_dict)
    writer = pd.ExcelWriter(filename, engine='xlsxwriter')
    combined_df.to_excel(writer, sheet_name='Monthly_Data')
    writer.save()

def save_csv(X, dest_path):
    X.to_csv(dest_path, index = False)

def dt2s(datetime_obj):
    time = datetime.time.strftime(datetime_obj,"%H:%M:%S.%f")
    seconds = int(time[:2])*60**2 + int(time[3:5])*60 + float(time[6:])
    return round(seconds,2)

def time_HHMM(in_string):
    return in_string[:in_string.find(':', 3)]
def time_MMSS(in_string):
    return in_string[in_string.find(':')+1:]

def hours_long_main():
    #Load excel file into 
    X = pd.DataFrame(pd.read_excel(historic))
    Y = pd.DataFrame(pd.read_excel(current))

    # Use this area to grab necessary info about the graph
    plant_names = list(set(X['Plant Name'].to_list()))


    # Use this area to create a plant dictionary, with each plant being a dict
    for p in plant_names: 
        plant_dict[p] = {}
        root_die_historic_dict[p] = {}
        current_month_dict[p] = []

    # go through year's data, groups Root Die with same ones in each plant
    for index, row in X.iterrows(): # check every row in 2021 year
        #throw in the row into the Correct Plant inside
        total_time = row['Setup']+row['Run']+row['Down'] 
        if row['Confirmed Qty']> 0 and total_time > 0:
            if row['Root Die'] not in plant_dict[row['Plant Name']]:
                plant_dict[row['Plant Name']][row['Root Die']] = []
            plant_dict[row['Plant Name']][row['Root Die']].append(list(row))
    


    # Now go through each plant, and create a historic die lookup table
    # matching Die Order to pairs of WC+Rate
    for plant in plant_dict:
        for root_die_list in plant_dict[plant].items(): 
            for root_die_instance in root_die_list:
                actual_rate = round(root_die_instance[4]/sum(root_die_instance[1:4]), 2)
                #each array will append the actual rate as meta data
                root_die_instance.append(actual_rate)

                #Append this pair (WC and Rate) into the dictionary
                wc_rate_pair = [root_die_instance[5], root_die_instance[8]]
                if root_die_instance[6] not in root_die_historic_dict[plant]:
                    root_die_historic_dict[plant][root_die_instance[6]] = []
                root_die_historic_dict[plant][root_die_instance[6]].append(wc_rate_pair)


    # Now that I have the Root Die list, we can go through the May data
    current_month = []
    for rowy in Y.iterrows():
        # print(list(rowy))
        row_time = (rowy[1]+rowy[2]+rowy[3])
        if row_time > 0 and rowy[4] > 0: # Only positive time and quantity allowed
            actual_rate = round(rowy[4]/row_time, 2)
            current_month.append(list(rowy))



    # In Montth's data, append Best Asset, Schedule, 
    for root_die in current_month:
        #['Lumberton', 1.0, 39.185, 7.035, 495000, 7413, 57231, '4/16 CAN NTF NITRO']
        root_string = str(root_die[6])
        #Calculate actual rate and append to row
        actual_rate = round(root_die[4]/sum(root_die[1:4]), 2)

        #Now go find the best asset rate and machine number
        max_set, repeat_set = [], []
        max_rate, repeat_rate= -10, -10

        if root_string in root_die_historic_dict[root_die[0]]: #current root die exists in the plant records
            for rate_pair in root_die_historic_dict[root_die[0]][root_string]: #check each rate pair
                if rate_pair[1] > max_rate: #If current pair > max rate, update the pair
                    max_rate = rate_pair[1]
                    max_set = rate_pair
                if rate_pair[0] == root_die[5] and rate_pair[1] > repeat_rate:
                    repeat_rate = rate_pair[1] 
                    repeat_set = rate_pair
        
        #Shove these values into current month as helper columns
        root_die.append(max_set[0])
        root_die.append(round(root_die[4]/max_set[1],2)) #append Best Asset Historic rate's hours
        root_die.append(round(root_die[4]/repeat_set[1],2)) #append Scheduled's Historic rate' hours
        root_die.append(round(root_die[4]/actual_rate, 2)) #append actual rate's hours


    # Each row goes into their correct Plant
    for cm_row in current_month:
        current_month_dict[cm_row[0]].append(cm_row)

    # Finally Time to Display the results!
    for plant, month_data in current_month_dict.items():
        opt_hours, schedule_hours, actual_hours = 0, 0, 0
        gamble_won, gamble_lost = 0, 0
        total_root_die = len(month_data)
        for rd_data in month_data:
            #['Newton', 3.0, 10.834, 5.016, 339500, 7404, 15646, 'NC FULLY COOKED BACON', 7401, 8.0, 15.39, 18.85]
            opt_hours += rd_data[9]
            schedule_hours += rd_data[10]
            actual_hours += rd_data[11]

            #If BA hours > Actual Hours
            if rd_data[5] != rd_data[8]:
                if rd_data[9] > rd_data[11]:
                    gamble_won += 1
                else:
                    gamble_lost += 1


        # FOR DISPLAYING THE RESULTS
        print("For Plant %s, the results are: " % plant)
        print(round(opt_hours, 2), round(schedule_hours, 2), round(actual_hours, 2))
        print("Out of %d root dies, Underperform: %d, Outperformed: %d" %
             (total_root_die, gamble_lost, gamble_won))
            
        print(gamble_lost, gamble_won, gamble_lost+gamble_won,total_root_die)
        print(gamble_won/(gamble_lost+gamble_won), gamble_lost/(gamble_lost+gamble_won), 
            (total_root_die-(gamble_lost+gamble_won))/total_root_die)

        print(" ")

def subtotal_month_plant_packager(plant_info):

    # take in all the rows of month-plant
    subtotal_array = []

    for i in range(len(plant_info[0])):
        print(i, plant_info[0][i], type(plant_info[0][i]))
        if i == 2: #grab the name of the plant
            subtotal_array.append(plant_info[0][i][0].upper() + plant_info[0][i][1:].lower())
        if i > 8 and i < 12: # Find the sum for these rows
            sum = 0
            for row in plant_info:
                print(row[i])
                sum += row[i]
            subtotal_array.append(round(sum, 2))
        elif i > 11: #Find the average for these
            row_count = 0
            plant_sum = 0
            for row in plant_info:
                print(row[i])
                row_count += 1
                print(type(row[i]))
                # if type(row[i]) == "<class 'datetime.time'>":
                if type(row[i]) == datetime.time:
                    seconds = datetime.time.strftime(row[i],"%H:%M:%S")
                    seconds = seconds.total_seconds()
                    print(seconds)



                #WORK ON THIS LATER
                #MAKE ALL DATE TIMES INTO SECONDS

                #THEN FIND THE AVERAGE OF SECONDS, TURN AVERAGE INTO DATE TIME
                # plant_sum += row[i]
            # subtotal_array.append(plant_sum/row_count)
        print('\n')
    print('End of the function, ')
    print(subtotal_array)


def graph1(bic, selected_plant, image_path):

    graph1_dict = {}
    month_num = 0
    for month in mr_month_dict:
        for plant in mr_month_dict[month]:
            curr_plant = plant.title() #current plant
            bic_plant = bic[2].title()
            if curr_plant == bic_plant: #if best plant found, add work center and months
                if bic_plant not in graph1_dict: #if BA plant not in dict, make a dict
                    graph1_dict[bic_plant] = {}
                if bic[3] not in graph1_dict[bic_plant]: #if WC not in dict, make array
                    graph1_dict[bic_plant][bic[3]] = []

                for wc_BA in mr_month_dict[month][plant]:
                    if wc_BA[3] == bic[3]: 
                        month_num+=1
                        avg_total_MR_time = dt2s(wc_BA[27]) + dt2s(wc_BA[30]) + dt2s(wc_BA[36])
                        # print('Work center is: ', wc_BA[3], month, plant, avg_total_MR_time)
                        graph1_dict[bic_plant][bic[3]].append(round(avg_total_MR_time/60,2))
            if curr_plant == selected_plant:
                if selected_plant not in graph1_dict:
                    graph1_dict[selected_plant] = {}
                for wc in mr_month_dict[month][plant]:
                    if wc[3] not in graph1_dict[selected_plant]:
                        graph1_dict[selected_plant][wc[3]] = []

                    avg_total_MR_time = dt2s(wc[27]) + dt2s(wc[30]) + dt2s(wc[36])
                    # print('Work center is: ', wc[3], month, plant, avg_total_MR_time)
                    graph1_dict[selected_plant][wc[3]].append(round(avg_total_MR_time/60, 2))

    month_tick_list = []
    month_arr_big = []
    for plant in graph1_dict:
        for wc in graph1_dict[plant]:
            month_tick_list.append(plant + " " + wc) #grab all the month tick list
    for i in range(month_num):
        m_arr = []
        for plant in graph1_dict:
            for wc in graph1_dict[plant]:
                m_arr.append(graph1_dict[plant][wc][i]) #grab the data month by month
        month_arr_big.append(m_arr)
    # print(month_tick_list, month_arr_big)

    bic_pos = -1
    for count, tick in enumerate(month_tick_list):
        if bic[2] in tick:
            month_tick_list[count] = "Best in Class"
            bic_pos = count

    # Keep only BIC All Time
    zero_list = []
    for zero in month_arr_big: zero_list.append(zero[bic_pos])
    for zero_1 in month_arr_big:
        if zero_1[bic_pos] != min(zero_list): zero_1[bic_pos] = 0

    #throw the arrays into the graph 1 tool
    bar1_grapher(month_tick_list, month_arr_big, image_path)
def bar1_grapher(X, vals, image_path, width=0.8):
    n = len(vals)

    _X = np.arange(len(X))
    plt.subplots(figsize=(10, 6)) #Future note, always fig size before plotting
    colors = color_list
    for i in range(n):
        plt.bar(_X - width/2. + i/float(n)*width, vals[i], 
                width=width/float(n), align="edge", color = colors[i])   
    
    plt.xlabel('Plant & Work Center', size = 12)
    plt.title('Average Total MR Time Month over Month', size = 15)
    plt.legend(month_label)

    # Y-Tick Factory
    largest_time = max(max(x) for x in vals) * 1.08
    y_tick_num = 9
    y_spaced = [largest_time * ((i+1)/y_tick_num) for i in range(y_tick_num)]
    y_timedelta = [str(timedelta(seconds=int(s))) for s in y_spaced]
    if largest_time >= 3600:
        y_timedelta_str = [time_HHMM(s) for s in y_timedelta]
        plt.ylabel('Time (h)', size = 12)
    else:
        y_timedelta_str = [time_MMSS(s) for s in y_timedelta]
        plt.ylabel('Time (m)', size = 12)
    plt.yticks(y_spaced, y_timedelta_str, size = 7)

    plt.xticks(_X, X, size=7)
    plt.grid(axis = 'y', linestyle = '--', linewidth = 0.3)
    
    if save_fig: plt.savefig('%s/%s' % (image_path, 'avg_total_mr_time_month_over_month_G1.png'))
    else: plt.show()
    plt.close()

def graph2(bic, selected_plant, image_path):

    bic_MR_time = dt2s(bic[27]) + dt2s(bic[30]) + dt2s(bic[36])

    graph1_dict = {}
    month_num = 0
    for month in mr_month_dict:
        for plant in mr_month_dict[month]:
            curr_plant = plant.title() #current plant
            bic_plant = bic[2].title()
            if curr_plant == bic_plant: #if best plant found, add work center and months
                if bic_plant not in graph1_dict: #if BA plant not in dict, make a dict
                    graph1_dict[bic_plant] = {}
                if bic[3] not in graph1_dict[bic_plant]: #if WC not in dict, make array
                    graph1_dict[bic_plant][bic[3]] = []

                for wc_BA in mr_month_dict[month][plant]:
                    if wc_BA[3] == bic[3]: 
                        month_num+=1
                        avg_total_MR_time = dt2s(wc_BA[27]) + dt2s(wc_BA[30]) + dt2s(wc_BA[36])
                        lost_sheets = ((avg_total_MR_time-bic_MR_time)/(60**2))*wc_BA[10]*wc_BA[20]
                        graph1_dict[bic_plant][bic[3]].append(lost_sheets)
            if curr_plant == selected_plant:
                if selected_plant not in graph1_dict:
                    graph1_dict[selected_plant] = {}
                for wc in mr_month_dict[month][plant]:
                    if wc[3] not in graph1_dict[selected_plant]:
                        graph1_dict[selected_plant][wc[3]] = []
                    avg_total_MR_time = dt2s(wc[27]) + dt2s(wc[30]) + dt2s(wc[36])
                    lost_sheets = ((avg_total_MR_time-bic_MR_time)/(60**2))*wc[10]*wc[20]
                    graph1_dict[selected_plant][wc[3]].append(lost_sheets)

    month_tick_list = []
    month_arr_big = []
    for plant in graph1_dict:
        for wc in graph1_dict[plant]:
            month_tick_list.append(plant + " " + wc) #grab all the month tick list
    for i in range(month_num):
        m_arr = []
        for plant in graph1_dict:
            for wc in graph1_dict[plant]:
                m_arr.append(graph1_dict[plant][wc][i]) #grab the data month by month
        month_arr_big.append(m_arr)
    # print(month_tick_list, month_arr_big)

    bic_pos = -1
    for count, tick in enumerate(month_tick_list):
        if bic[2] in tick:
            month_tick_list[count] = "Best in Class"
            bic_pos = count

    del month_tick_list[bic_pos]
    for m in month_arr_big:
        del m[bic_pos]

    #throw the arrays into the graph 1 tool
    bar2_grapher(month_tick_list, month_arr_big, image_path)
def bar2_grapher(X, vals, image_path, width=0.8):
    n = len(vals)
    _X = np.arange(len(X))
    plt.subplots(figsize=(10, 6)) #Future note, always fig size before plotting
    colors = color_list
    for i in range(n):
        plt.bar(_X - width/2. + i/float(n)*width, vals[i], 
                width=width/float(n), align="edge", color = colors[i])   
    
    plt.title('Make Ready Lost Sheets vs Best in Class Month over Month', size = 15)
    plt.legend(month_label)
    plt.xlabel('Plant & Work Center', size = 12)
    largest_num = max(max(x) for x in vals)
    if largest_num > 10**6: plt.ylabel('Millions', size = 12)
    plt.xticks(_X, X, size=7)
    plt.grid(axis = 'y', linestyle = '--', linewidth = 0.3)
    
    if save_fig: plt.savefig('%s/%s' % (image_path, 'Mmr_lost_sheets_vs_bic_mom_G2.png'))
    else: plt.show()
    plt.close()

def graph3(bic, selected_plant, image_path):

    wc_row = []
    ana_month = ''
    # Grab the latest month
    for month in mr_month_dict:
        ana_month = month
    for month in mr_month_dict:
        for plant in mr_month_dict[month]:
            curr_plant = plant.title()
            bic_plant = bic[2].title()
            if curr_plant == bic_plant:
                for wc_bic in mr_month_dict[month][plant]:
                    if bic[3] == wc_bic[3] and month == str(bic[5])[:10]:
                        sum_of_ss_perc = round(100*(dt2s(wc_bic[38])+dt2s(wc_bic[32])+dt2s(wc_bic[29]))/(dt2s(wc_bic[15])+dt2s(wc_bic[27])+dt2s(wc_bic[30])), 2)
                        wc_row.append([curr_plant+ " "+ wc_bic[3], dt2s(wc_bic[15]), dt2s(wc_bic[27]), dt2s(wc_bic[30]), sum_of_ss_perc])
            if curr_plant == selected_plant and month == ana_month: #in the selected plant
                for wc in mr_month_dict[month][plant]:
                    sum_of_ss_perc = round(100*(dt2s(wc[38])+dt2s(wc[32])+dt2s(wc[29]))/(dt2s(wc[15])+dt2s(wc[27])+dt2s(wc[30])), 2)
                    wc_row.append([curr_plant+ " "+ wc[3], dt2s(wc[15]), dt2s(wc[27]), dt2s(wc[30]), sum_of_ss_perc])

    presentation_row = []
    for row in wc_row:
        if row[0] not in presentation_row:
            presentation_row.append(row[0])

    #Intended for this to combine multiple months, but still works
    #for a single month, so no changes needed
    
    tick_list_serial_row = []
    
    for serial in presentation_row:
        serial_list = []
        for row in wc_row:
            if serial == row[0]:
                serial_list.append(row)
        combined_serial_list = []
        combined_serial_list.append(serial_list[0][0]) #grab the head
        ot_sum, basic_MR, ft_time, ss_time_perc = 0, 0, 0, 0
        for ser in serial_list:
            ot_sum += ser[1]
            basic_MR += ser[2]
            ft_time += ser[3]
            ss_time_perc += ser[4]
        combined_serial_list.append(ot_sum)
        combined_serial_list.append(basic_MR)
        combined_serial_list.append(ft_time)
        combined_serial_list.append(ss_time_perc)
        tick_list_serial_row.append(combined_serial_list)
    
    #Now we can finally populate the three arrays
    ot_arr = []
    basic_MR_arr = []
    ft_time_arr = []
    ss_time_perc_arr = []
    for tick_list in tick_list_serial_row:
        ot_arr.append(tick_list[1])
        basic_MR_arr.append(tick_list[2])
        ft_time_arr.append(tick_list[3])
        ss_time_perc_arr.append(tick_list[4]) #percentage add
    tick_array = [ot_arr, basic_MR_arr, ft_time_arr, ss_time_perc_arr]

    for count, tick in enumerate(presentation_row):
        if bic[2] in tick:
            presentation_row[count] = "Best in Class"
    #Now I can finally plot all these
    bar3_grapher(presentation_row, tick_array, image_path)
def bar3_grapher(labels, avg_arr, image_path):

    width = 0.35       # the width of the bars: can also be len(x) sequence
    fig, ax = plt.subplots(figsize=(10, 6))
    colors = color_list
    ax2 = ax.twinx()
    ax.bar(labels, avg_arr[0], width, label='Sum of AVG Other Time', color = colors[0])
    ax.bar(labels, avg_arr[1], width, bottom=avg_arr[0], 
            label='Sum of AVG Basic MR Time', color = colors[1])
    ax.bar(labels, avg_arr[2], width, bottom=list(np.add(avg_arr[0], avg_arr[1]))
            , label='Sum of AVG Fine Tuning Time', color = colors[2])

    plt.plot(labels, avg_arr[3], linestyle='None', marker='D', color='b', label='Sum of SS Total %')
    
    # Y-Tick Factory
    vals = [0 for _ in range(len(avg_arr[0]))] #vals = #of WCs
    for i in range(len(avg_arr[0])): #for WC1
        for metric in avg_arr: #check each metric
            vals[i] += metric[i] #add each metric for wc1
    largest_time = max(vals) * 1.1

    y_tick_num = 8
    y_spaced = [largest_time * ((i+1)/y_tick_num) for i in range(y_tick_num)]
    y_timedelta = [str(timedelta(seconds=int(s))) for s in y_spaced]
    if largest_time >= 3600:
        y_timedelta_str = [time_HHMM(s) for s in y_timedelta]
        ax.set_ylabel('Time (h)', size = 12)
    else:
        y_timedelta_str = [time_MMSS(s) for s in y_timedelta]
        ax.set_ylabel('Time (m)', size = 12)
    #Note! Since double plots, you need set_yticks not y_ticks
    ax.set_yticks(y_spaced, y_timedelta_str, size = 7)

    ax2.set_ylabel('Sum of SS %', size = 12) 
    ax.set_title('Average Total Make Ready Time', size = 15)
    ax2.set_ylim([0, 100])

    # Put the legends in correct places
    ax.legend(loc='upper right')
    plt.legend(loc='upper left')
    ax.grid(axis = 'y', linestyle = '--', linewidth = 0.3)
    
    if save_fig: plt.savefig('%s/%s' % (image_path, 'avg_total_mr_time_G3.png'))
    else: plt.show() 
    plt.close()

def graph4(bic, selected_plant, image_path):
    wc_row = []
    bic_MR_time = dt2s(bic[27])
    bic_FT_time = dt2s(bic[30])
    bic_OT_time = dt2s(bic[36])

    ana_month = ''
    # Grab the latest month
    for month in mr_month_dict:
        ana_month = month
    for month in mr_month_dict:
        for plant in mr_month_dict[month]:
            curr_plant = plant.title()
            bic_plant = bic[2].title()
            if curr_plant == bic_plant and month == ana_month:
                for wc_bic in mr_month_dict[month][plant]:
                    if bic[3] == wc_bic[3]:
                        MR = wc_bic[10] * wc_bic[20] * (dt2s(wc_bic[27]) - bic_MR_time)/(60**2)
                        FT = wc_bic[10] * wc_bic[20] * (dt2s(wc_bic[30]) - bic_FT_time)/(60**2)
                        OT = wc_bic[10] * wc_bic[20] * (dt2s(wc_bic[36]) - bic_OT_time)/(60**2)
                        wc_row.append([curr_plant+ " "+ wc_bic[3], OT, MR, FT])
                        
            if curr_plant == selected_plant and month == ana_month: #in the selected plant
                for wc in mr_month_dict[month][plant]:

                    MR = wc[10] * wc[20] * (dt2s(wc[27]) - bic_MR_time)/(60**2)
                    FT = wc[10] * wc[20] * (dt2s(wc[30]) - bic_FT_time)/(60**2)
                    OT = wc[10] * wc[20] * (dt2s(wc[36]) - bic_OT_time)/(60**2)
                    wc_row.append([curr_plant+ " "+ wc[3], OT, MR, FT])


    
    presentation_row = []
    for row in wc_row:
        if row[0] not in presentation_row:
            presentation_row.append(row[0])

    #Only use this if I want all months
    tick_list_serial_row = []
    for serial in presentation_row:
        serial_list = []
        for row in wc_row:
            if serial == row[0]:
                # print(row)
                serial_list.append(row)

        combined_serial_list = []
        combined_serial_list.append(serial_list[0][0]) #grab the head
        ot_sum, basic_MR, ft_time = 0, 0, 0
        for ser in serial_list:
            ot_sum += ser[1]
            basic_MR += ser[2]
            ft_time += ser[3]
        combined_serial_list.append(ot_sum)
        combined_serial_list.append(basic_MR)
        combined_serial_list.append(ft_time)
        tick_list_serial_row.append(combined_serial_list)
    
    #Now we can finally populate the three arrays
    ot_arr = []
    basic_MR_arr = []
    ft_time_arr = []
    for tick_list in tick_list_serial_row:
        ot_arr.append(tick_list[1])
        basic_MR_arr.append(tick_list[2])
        ft_time_arr.append(tick_list[3])
    tick_array = [ot_arr, basic_MR_arr, ft_time_arr]

    #Now I can finally plot all these
    bic_pos = -1
    for count, tick in enumerate(presentation_row):
        if bic[2] in tick:
            presentation_row[count] = "Best in Class"
            bic_pos = count
    del presentation_row[bic_pos]
    for m in tick_array:
        del m[bic_pos]
    bar4_grapher(presentation_row, tick_array, image_path)
def bar4_grapher(labels, avg_arr, image_path):
    # print(avg_arr)
    width = 0.35       # the width of the bars: can also be len(x) sequence
    fig, ax = plt.subplots(figsize=(10, 6))
    colors = color_list
    ax.bar(labels, avg_arr[0], width, label='Sum of Other Time Lost Sheets', color = colors[0])
    ax.bar(labels, avg_arr[1], width, bottom=avg_arr[0], 
            label='Sum of Basic MR Time Lost Sheets', color = colors[1])
    ax.bar(labels, avg_arr[2], width, bottom=list(np.add(avg_arr[0], avg_arr[1]))
            , label='Sum of Fine Tune Time Lost Sheets', color = colors[2])


    largest_num = max(max(x) for x in avg_arr)
    if largest_num > 10**6: plt.ylabel('Millions', size = 12)
    ax.set_title('Total MR Lost Sheets vs Best in Class', size = 15)
    ax.legend(loc='upper right')
    plt.grid(axis = 'y', linestyle = '--', linewidth = 0.3)
    if save_fig: plt.savefig('%s/%s' % (image_path, 'mr_lost_sheets_G4.png'))
    else: plt.show() 
    plt.close()

def graph5(bic, selected_plant, image_path):
    wc_row = []
    ana_month = ''
    # Grab the latest month
    for month in mr_month_dict:
        ana_month = month
    for month in mr_month_dict:
        for plant in mr_month_dict[month]:
            curr_plant = plant.title()
            bic_plant = bic[2].title()
            if curr_plant == bic_plant and month == str(bic[5])[:10]:
                for wc_bic in mr_month_dict[month][plant]:
                    if bic[3] == wc_bic[3]:
                        wc_row.append([curr_plant+ " "+ wc_bic[3], dt2s(wc_bic[15]), dt2s(wc_bic[38]), 100*dt2s(wc_bic[38])/dt2s(wc_bic[15])])
            if curr_plant == selected_plant and month == ana_month: #in the selected plant
                for wc in mr_month_dict[month][plant]:
                    wc_row.append([curr_plant+ " "+ wc[3], dt2s(wc[15]), dt2s(wc[38]), 100*dt2s(wc[38])/dt2s(wc[15])])
    
    presentation_row = []
    for row in wc_row:
        if row[0] not in presentation_row:
            presentation_row.append(row[0])

    #Only use this if I want all months
    tick_list_serial_row = []
    for serial in presentation_row:
        serial_list = []
        for row in wc_row:
            if serial == row[0]:
                # print(row)
                serial_list.append(row)

        combined_serial_list = []
        combined_serial_list.append(serial_list[0][0]) #grab the head
        ot_sum, basic_MR, ft_time = 0, 0, 0
        for ser in serial_list:
            ot_sum += ser[1]
            basic_MR += ser[2]
            ft_time += ser[3]
        combined_serial_list.append(ot_sum)
        combined_serial_list.append(basic_MR)
        combined_serial_list.append(ft_time)
        tick_list_serial_row.append(combined_serial_list)
    
    #Now we can finally populate the three arrays
    ot_arr = []
    basic_MR_arr = []
    ft_time_arr = []
    for tick_list in tick_list_serial_row:
        ot_arr.append(tick_list[1])
        basic_MR_arr.append(tick_list[2])
        ft_time_arr.append(tick_list[3])
    tick_array = [ot_arr, basic_MR_arr, ft_time_arr]

    #Now I can finally plot all these
    for count, tick in enumerate(presentation_row):
        if bic[2] in tick:
            presentation_row[count] = "Best in Class"
    bar5_grapher(presentation_row, tick_array, image_path)
def bar5_grapher(labels, avg_arr, image_path, width=0.8):
    side_arr = avg_arr[:2]
    n = len(side_arr)
    _X = np.arange(len(labels))
    fig, ax = plt.subplots(figsize=(10, 6)) #Future note, always fig size before plotting
    colors = color_list
    for i in range(n):
        plt.bar(_X - width/2. + i/float(n)*width, avg_arr[i], 
                width=width/float(n), align="edge", color = colors[i])   
    ax2 = ax.twinx()
    plt.plot(labels, avg_arr[2], linestyle='None', marker='D', color='b', label='Sum of SS Other %')
    # Y-Tick Factory
    largest_time = max(max(x) for x in side_arr) * 1.1
    y_tick_num = 8
    y_spaced = [largest_time * ((i+1)/y_tick_num) for i in range(y_tick_num)]
    y_timedelta = [str(timedelta(seconds=int(s))) for s in y_spaced]
    if largest_time >= 3600:
        y_timedelta_str = [time_HHMM(s) for s in y_timedelta]
        ax.set_ylabel('Time (h)', size = 12)
    else:
        y_timedelta_str = [time_MMSS(s) for s in y_timedelta]
        ax.set_ylabel('Time (m)', size = 12)
    #Note! Since double plots, you need set_yticks not y_ticks
    ax.set_yticks(y_spaced, y_timedelta_str, size = 7)
    ax2.set_ylabel('Sum of SS Other %', size = 12) 
    ax.set_title('Average Other Make Ready Time', size = 15)
    ax2.set_ylim([0, 100])

    # Put the legends in correct places
    ax.legend(['Sum of AVG Other Time', 'Sum of AVG SS Other Time'], loc='upper right')
    plt.legend(loc='upper left')
    ax.grid(axis = 'y', linestyle = '--', linewidth = 0.3)
    
    if save_fig: plt.savefig('%s/%s' % (image_path, 'avg_other_mr_time_G5.png'))
    else: plt.show() 
    plt.close()

def graph6(bic, selected_plant, image_path):
    wc_row = []
    bic_OT_time = dt2s(bic[36])
    ana_month = ''
    # Grab the latest month
    for month in mr_month_dict:
        ana_month = month
    for month in mr_month_dict:
        for plant in mr_month_dict[month]:
            curr_plant = plant.title()
            bic_plant = bic[2].title()
            if curr_plant == bic_plant and month == ana_month:
                for wc_bic in mr_month_dict[month][plant]:
                    if bic[3] == wc_bic[3]:
                        wc_row.append([curr_plant+ " "+ wc_bic[3], wc_bic[10] * wc_bic[20] * (dt2s(wc_bic[36]) - bic_OT_time)/(60**2)])
            if curr_plant == selected_plant and month == ana_month: #in the selected plant
                for wc in mr_month_dict[month][plant]:
                    wc_row.append([curr_plant+ " "+ wc[3], wc[10] * wc[20] * (dt2s(wc[36]) - bic_OT_time)/(60**2)])
    
    presentation_row = []
    for row in wc_row:
        if row[0] not in presentation_row:
            presentation_row.append(row[0])

    #Only use this if I want all months
    tick_list_serial_row = []
    for serial in presentation_row:
        serial_list = []
        for row in wc_row:
            if serial == row[0]:
                # print(row)
                serial_list.append(row)

        combined_serial_list = []
        combined_serial_list.append(serial_list[0][0]) #grab the head
        ot_sum = 0
        for ser in serial_list:
            ot_sum += ser[1]

        combined_serial_list.append(ot_sum)
        tick_list_serial_row.append(combined_serial_list)
    
    #Now we can finally populate the three arrays
    ot_arr = []
    for tick_list in tick_list_serial_row:
        ot_arr.append(tick_list[1])
    tick_array = [ot_arr]

    #Now I can finally plot all these
    bic_pos = -1
    for count, tick in enumerate(presentation_row):
        if bic[2] in tick:
            presentation_row[count] = "Best in Class"
            bic_pos = count
    del presentation_row[bic_pos]
    for m in tick_array:
        del m[bic_pos]
    bar6_grapher(presentation_row, tick_array, image_path)
def bar6_grapher(labels, avg_arr, image_path):
    width = 0.35       # the width of the bars: can also be len(x) sequence
    fig, ax = plt.subplots(figsize=(10, 6))
    colors = color_list
    ax.bar(labels, avg_arr[0], width, label='Average Blanket Wash Time', color = colors[0])
    
    ax.set_title('Other MR Lost Sheets vs Best in Class', size = 15)
    # ax.legend(['Sum of Other Time Lost Sheets'], loc='upper right')
    largest_num = max(max(x) for x in avg_arr)
    if largest_num > 10**6: ax.set_ylabel('Millions', size = 12)
    plt.grid(axis = 'y', linestyle = '--', linewidth = 0.3)
    if save_fig: plt.savefig('%s/%s' % (image_path, 'other_mr_lost_sheets_G6.png'))
    else: plt.show() 
    plt.close()

def graph7(bic, selected_plant, image_path):
    wc_row = []
    ana_month = ''
    # Grab the latest month
    for month in mr_month_dict:
        ana_month = month
    for month in mr_month_dict:
        for plant in mr_month_dict[month]:
            curr_plant = plant.title()
            bic_plant = bic[2].title()
            if curr_plant == bic_plant and month == str(bic[5])[:10]:
                for wc_bic in mr_month_dict[month][plant]:
                    if bic[3] == wc_bic[3]:
                        wc_row.append([curr_plant+ " "+ wc_bic[3], dt2s(wc_bic[27]), dt2s(wc_bic[29]), 100*dt2s(wc_bic[29])/dt2s(wc_bic[27])])
            if curr_plant == selected_plant and month == ana_month: #in the selected plant
                for wc in mr_month_dict[month][plant]:
                    wc_row.append([curr_plant+ " "+ wc[3], dt2s(wc[27]), dt2s(wc[29]), 100*dt2s(wc[29])/dt2s(wc[27])])
    
    presentation_row = []
    for row in wc_row:
        if row[0] not in presentation_row:
            presentation_row.append(row[0])

    #Only use this if I want all months
    tick_list_serial_row = []
    for serial in presentation_row:
        serial_list = []
        for row in wc_row:
            if serial == row[0]:
                # print(row)
                serial_list.append(row)

        combined_serial_list = []
        combined_serial_list.append(serial_list[0][0]) #grab the head
        ot_sum, basic_MR, ft_time = 0, 0, 0
        for ser in serial_list:
            ot_sum += ser[1]
            basic_MR += ser[2]
            ft_time += ser[3]
        combined_serial_list.append(ot_sum)
        combined_serial_list.append(basic_MR)
        combined_serial_list.append(ft_time)
        tick_list_serial_row.append(combined_serial_list)
    
    #Now we can finally populate the three arrays
    ot_arr = []
    basic_MR_arr = []
    ft_time_arr = []
    for tick_list in tick_list_serial_row:
        ot_arr.append(tick_list[1])
        basic_MR_arr.append(tick_list[2])
        ft_time_arr.append(tick_list[3])
    tick_array = [ot_arr, basic_MR_arr, ft_time_arr]

    #Now I can finally plot all these
    for count, tick in enumerate(presentation_row):
        if bic[2] in tick:
            presentation_row[count] = "Best in Class"
    bar7_grapher(presentation_row, tick_array, image_path)
def bar7_grapher(labels, avg_arr, image_path, width=0.8):
    side_arr = avg_arr[:2]
    n = len(side_arr)

    # Start graphing the side by side graphs
    _X = np.arange(len(labels))
    fig, ax = plt.subplots(figsize=(10, 6)) #Future note, always fig size before plotting
    colors = color_list
    for i in range(n):
        plt.bar(_X - width/2. + i/float(n)*width, avg_arr[i], 
                width=width/float(n), align="edge", color = colors[i])   

    #Created second axis, plot the %s
    ax2 = ax.twinx()
    plt.plot(labels, avg_arr[2], linestyle='None', marker='D', color='b', label='Sum of SS Basic %')
    
    # Y-Tick Factory
    largest_time = max(max(x) for x in side_arr) * 1.1
    y_tick_num = 8
    y_spaced = [largest_time * ((i+1)/y_tick_num) for i in range(y_tick_num)]
    y_timedelta = [str(timedelta(seconds=int(s))) for s in y_spaced]
    if largest_time >= 3600:
        y_timedelta_str = [time_HHMM(s) for s in y_timedelta]
        ax.set_ylabel('Time (h)', size = 12)
    else:
        y_timedelta_str = [time_MMSS(s) for s in y_timedelta]
        ax.set_ylabel('Time (m)', size = 12)
    #Note! Since double plots, you need set_yticks not y_ticks
    ax.set_yticks(y_spaced, y_timedelta_str, size = 7)

    ax2.set_ylabel('Sum of SS Basic %', size = 12) 
    ax2.set_ylim([0, 100])
    ax.set_title('Average Basic Make Ready Time', size = 15)
    ax.legend(['Sum of AVG. Basic MR Time', 'Sum of AVG. SS Basic Make Ready Time'], loc='upper right')
    plt.legend(loc='upper left')
    ax.grid(axis = 'y', linestyle = '--', linewidth = 0.3)
    
    if save_fig: plt.savefig('%s/%s' % (image_path, 'avg_basic_mr_time_G7.png'))
    else: plt.show() 
    plt.close()

def graph8(bic, selected_plant, image_path):
    wc_row = []
    bic_MR_time = dt2s(bic[27])
    ana_month = ''
    # Grab the latest month
    for month in mr_month_dict:
        ana_month = month
    for month in mr_month_dict:
        for plant in mr_month_dict[month]:
            curr_plant = plant.title()
            bic_plant = bic[2].title()
            if curr_plant == bic_plant and month == ana_month:
                for wc_bic in mr_month_dict[month][plant]:
                    if bic[3] == wc_bic[3]:
                        wc_row.append([curr_plant+ " "+ wc_bic[3], wc_bic[10] * wc_bic[20] * (dt2s(wc_bic[27]) - bic_MR_time)/(60**2)])
            if curr_plant == selected_plant and month == ana_month: #in the selected plant
                for wc in mr_month_dict[month][plant]:
                    wc_row.append([curr_plant+ " "+ wc[3], wc[10] * wc[20] * (dt2s(wc[27]) - bic_MR_time)/(60**2)])
    
    presentation_row = []
    for row in wc_row:
        if row[0] not in presentation_row:
            presentation_row.append(row[0])

    #Only use this if I want all months
    tick_list_serial_row = []
    for serial in presentation_row:
        serial_list = []
        for row in wc_row:
            if serial == row[0]:
                # print(row)
                serial_list.append(row)

        combined_serial_list = []
        combined_serial_list.append(serial_list[0][0]) #grab the head
        ot_sum = 0
        for ser in serial_list:
            ot_sum += ser[1]

        combined_serial_list.append(ot_sum)
        tick_list_serial_row.append(combined_serial_list)
    
    #Now we can finally populate the three arrays
    ot_arr = []

    for tick_list in tick_list_serial_row:
        ot_arr.append(tick_list[1])

    tick_array = [ot_arr]

    #Now I can finally plot all these
    bic_pos = -1
    for count, tick in enumerate(presentation_row):
        if bic[2] in tick:
            presentation_row[count] = "Best in Class"
            bic_pos = count
    del presentation_row[bic_pos]
    for m in tick_array:
        del m[bic_pos]
    bar8_grapher(presentation_row, tick_array, image_path)
def bar8_grapher(labels, avg_arr, image_path):
    width = 0.35       # the width of the bars: can also be len(x) sequence
    fig, ax = plt.subplots(figsize=(10, 6))
    colors = color_list
    ax.bar(labels, avg_arr[0], width, label='Average Blanket Wash Time', color = colors[0])
    
    ax.set_title('Other MR Lost Sheets vs Best in Class', size = 15)
    # ax.legend(['Sum of Basic Make Ready Time Lost Sheets'], loc='upper right')
    largest_num = max(max(x) for x in avg_arr)
    if largest_num > 10**6: ax.set_ylabel('Millions', size = 12)
    plt.grid(axis = 'y', linestyle = '--', linewidth = 0.3)
    if save_fig: plt.savefig('%s/%s' % (image_path, 'other_mr_lost_sheet_G8.png'))
    else: plt.show() 
    plt.close()

def graph9(bic, selected_plant, image_path):
    wc_row = []
    ana_month = ''
    # Grab the latest month
    for month in mr_month_dict:
        ana_month = month
    for month in mr_month_dict:
        for plant in mr_month_dict[month]:
            curr_plant = plant.title()
            bic_plant = bic[2].title()
            if curr_plant == bic_plant and month == str(bic[5])[:10]:
                for wc_bic in mr_month_dict[month][plant]:
                    if bic[3] == wc_bic[3]:
                        wc_row.append([curr_plant+ " "+ wc_bic[3], dt2s(wc_bic[30]), dt2s(wc_bic[32]), 100*dt2s(wc_bic[32])/dt2s(wc_bic[30])])
            if curr_plant == selected_plant and month == ana_month: #in the selected plant
                for wc in mr_month_dict[month][plant]:
                    wc_row.append([curr_plant+ " "+ wc[3], dt2s(wc[30]), dt2s(wc[32]), 100*dt2s(wc[32])/dt2s(wc[30])])
    
    presentation_row = []
    for row in wc_row:
        if row[0] not in presentation_row:
            presentation_row.append(row[0])

    #Only use this if I want all months
    tick_list_serial_row = []
    for serial in presentation_row:
        serial_list = []
        for row in wc_row:
            if serial == row[0]:
                # print(row)
                serial_list.append(row)

        combined_serial_list = []
        combined_serial_list.append(serial_list[0][0]) #grab the head
        ot_sum, basic_MR, ft_time = 0, 0, 0
        for ser in serial_list:
            ot_sum += ser[1]
            basic_MR += ser[2]
            ft_time += ser[3]
        combined_serial_list.append(ot_sum)
        combined_serial_list.append(basic_MR)
        combined_serial_list.append(ft_time)
        tick_list_serial_row.append(combined_serial_list)
    
    #Now we can finally populate the three arrays
    ot_arr = []
    basic_MR_arr = []
    ft_time_arr = []
    for tick_list in tick_list_serial_row:
        ot_arr.append(tick_list[1])
        basic_MR_arr.append(tick_list[2])
        ft_time_arr.append(tick_list[3])
    tick_array = [ot_arr, basic_MR_arr, ft_time_arr]

    #Now I can finally plot all these
    for count, tick in enumerate(presentation_row):
        if bic[2] in tick:
            presentation_row[count] = "Best in Class"
    bar9_grapher(presentation_row, tick_array, image_path)
def bar9_grapher(labels, avg_arr, image_path, width=0.8):

    side_arr = avg_arr[:2]
    n = len(side_arr)

    # Start graphing the side by side graphs
    _X = np.arange(len(labels))
    fig, ax = plt.subplots(figsize=(10, 6)) #Future note, always fig size before plotting
    colors = color_list
    for i in range(n):
        plt.bar(_X - width/2. + i/float(n)*width, avg_arr[i], 
                width=width/float(n), align="edge", color = colors[i])   

    #Created second axis, plot the %s
    ax2 = ax.twinx()
    plt.plot(labels, avg_arr[2], linestyle='None', marker='D', color='b', label='Sum of SS Fine %')
    
    # Y-Tick Factory
    largest_time = max(max(x) for x in side_arr) * 1.1
    y_tick_num = 8
    y_spaced = [largest_time * ((i+1)/y_tick_num) for i in range(y_tick_num)]
    y_timedelta = [str(timedelta(seconds=int(s))) for s in y_spaced]
    if largest_time >= 3600:
        y_timedelta_str = [time_HHMM(s) for s in y_timedelta]
        ax.set_ylabel('Time (h)', size = 12)
    else:
        y_timedelta_str = [time_MMSS(s) for s in y_timedelta]
        ax.set_ylabel('Time (m)', size = 12)
    #Note! Since double plots, you need set_yticks not y_ticks
    ax.set_yticks(y_spaced, y_timedelta_str, size = 7)

    ax2.set_ylabel('Sum of SS Fine %', size = 12) 
    ax2.set_ylim([0, 100])
    ax.set_title('Average Fine Tuning Make Ready Time', size = 15)
    ax.legend(['Sum of AVG. Fine Tuning Time', 'Sum of AVG. SS Fine Tuning Time'], loc='upper right')
    plt.legend(loc='upper left')
    ax.grid(axis = 'y', linestyle = '--', linewidth = 0.3)
    
    if save_fig: plt.savefig('%s/%s' % (image_path, 'avg_ft_mr_time_G9.png'))
    else: plt.show() 
    plt.close()

def graph10(bic, selected_plant, image_path):
    wc_row = []
    bic_FT_time = dt2s(bic[30])
    ana_month = ''
    # Grab the latest month
    for month in mr_month_dict:
        ana_month = month
    for month in mr_month_dict:
        for plant in mr_month_dict[month]:
            curr_plant = plant.title()
            bic_plant = bic[2].title()
            if curr_plant == bic_plant and month == ana_month:
                for wc_bic in mr_month_dict[month][plant]:
                    if bic[3] == wc_bic[3]:
                        wc_row.append([curr_plant+ " "+ wc_bic[3], wc_bic[10] * wc_bic[20] * (dt2s(wc_bic[30]) - bic_FT_time)/(60**2)])
            if curr_plant == selected_plant and month == ana_month: #in the selected plant
                for wc in mr_month_dict[month][plant]:
                    wc_row.append([curr_plant+ " "+ wc[3], wc[10] * wc[20] * (dt2s(wc[30]) - bic_FT_time)/(60**2)])
    
    presentation_row = []
    for row in wc_row:
        if row[0] not in presentation_row:
            presentation_row.append(row[0])

    #Only use this if I want all months
    tick_list_serial_row = []
    for serial in presentation_row:
        serial_list = []
        for row in wc_row:
            if serial == row[0]:
                # print(row)
                serial_list.append(row)

        combined_serial_list = []
        combined_serial_list.append(serial_list[0][0]) #grab the head
        ot_sum = 0
        for ser in serial_list:
            ot_sum += ser[1]

        combined_serial_list.append(ot_sum)
        tick_list_serial_row.append(combined_serial_list)
    
    #Now we can finally populate the three arrays
    ot_arr = []

    for tick_list in tick_list_serial_row:
        ot_arr.append(tick_list[1])

    tick_array = [ot_arr]

    #Now I can finally plot all these
    bic_pos = -1
    for count, tick in enumerate(presentation_row):
        if bic[2] in tick:
            presentation_row[count] = "Best in Class"
            bic_pos = count
    del presentation_row[bic_pos]
    for m in tick_array:
        del m[bic_pos]
    bar10_grapher(presentation_row, tick_array, image_path)
def bar10_grapher(labels, avg_arr, image_path):
    width = 0.35       # the width of the bars: can also be len(x) sequence
    fig, ax = plt.subplots(figsize=(10, 6))
    colors = color_list
    ax.bar(labels, avg_arr[0], width, label='Average Blanket Wash Time', color = colors[0])
    
    ax.set_title('Other MR Lost Sheets vs Best in Class', size = 15)
    # ax.legend(['Sum of Fine Tune Time Lost Sheets'], loc='upper right')
    largest_num = max(max(x) for x in avg_arr)
    if largest_num > 10**6: ax.set_ylabel('Millions', size = 12)
    plt.grid(axis = 'y', linestyle = '--', linewidth = 0.3)
    if save_fig: plt.savefig('%s/%s' % (image_path, 'other_mr_lost_sheets_G10.png'))
    else: plt.show() 
    plt.close()

def graph11(bic, selected_plant, image_path):
    wc_row = []
    ana_month = ''
    # Grab the latest month
    for month in mr_month_dict:
        ana_month = month
    for month in mr_month_dict:
        for plant in mr_month_dict[month]:
            curr_plant = plant.title()
            bic_plant = bic[2].title()
            if curr_plant == bic_plant and month == str(bic[5])[:10]:
                for wc_bic in mr_month_dict[month][plant]:
                    if bic[3] == wc_bic[3]:
                        wc_row.append([curr_plant+ " "+ wc_bic[3], dt2s(wc_bic[38]), dt2s(wc_bic[29]), dt2s(wc_bic[32])])
            if curr_plant == selected_plant and month == ana_month: #in the selected plant
                for wc in mr_month_dict[month][plant]:
                    wc_row.append([curr_plant+ " "+ wc[3], dt2s(wc[38]), dt2s(wc[29]), dt2s(wc[32])])
    
    presentation_row = []
    for row in wc_row:
        if row[0] not in presentation_row:
            presentation_row.append(row[0])

    #Only use this if I want all months
    tick_list_serial_row = []
    for serial in presentation_row:
        serial_list = []
        for row in wc_row:
            if serial == row[0]:
                # print(row)
                serial_list.append(row)

        combined_serial_list = []
        combined_serial_list.append(serial_list[0][0]) #grab the head
        ot_sum, basic_MR, ft_time = 0, 0, 0
        for ser in serial_list:
            ot_sum += ser[1]
            basic_MR += ser[2]
            ft_time += ser[3]
        combined_serial_list.append(ot_sum)
        combined_serial_list.append(basic_MR)
        combined_serial_list.append(ft_time)
        tick_list_serial_row.append(combined_serial_list)
    
    #Now we can finally populate the three arrays
    ot_arr = []
    basic_MR_arr = []
    ft_time_arr = []
    for tick_list in tick_list_serial_row:
        ot_arr.append(tick_list[1])
        basic_MR_arr.append(tick_list[2])
        ft_time_arr.append(tick_list[3])
    tick_array = [ot_arr, basic_MR_arr, ft_time_arr]

    #Now I can finally plot all these
    for count, tick in enumerate(presentation_row):
        if bic[2] in tick:
            presentation_row[count] = "Best in Class"
    bar11_grapher(presentation_row, tick_array, image_path)
def bar11_grapher(labels, avg_arr, image_path):
    width = 0.35       # the width of the bars: can also be len(x) sequence
    fig, ax = plt.subplots(figsize=(10, 6))
    colors = color_list
    ax.bar(labels, avg_arr[0], width, label='Sum of AVG SS Other Time', color = colors[0])
    ax.bar(labels, avg_arr[1], width, bottom=avg_arr[0], 
            label='Sum of AVG SS Basic Make Ready Time', color = colors[1])
    ax.bar(labels, avg_arr[2], width, bottom=list(np.add(avg_arr[0], avg_arr[1]))
            , label='Sum of AVG SS Fine Tuning Time', color = colors[2])

    # Y-Tick Factory
    vals = [0 for _ in range(len(avg_arr[0]))]
    for i in range(len(avg_arr[0])): #for WC1
        for metric in avg_arr: 
            vals[i] += metric[i]
    largest_time = max(vals) * 1.1
    y_tick_num = 9
    y_spaced = [largest_time * ((i+1)/y_tick_num) for i in range(y_tick_num)]
    y_timedelta = [str(timedelta(seconds=int(s))) for s in y_spaced]
    if largest_time >= 3600:
        y_timedelta_str = [time_HHMM(s) for s in y_timedelta]
        ax.set_ylabel('Time (h)', size = 12)
    else:
        y_timedelta_str = [time_MMSS(s) for s in y_timedelta]
        ax.set_ylabel('Time (m)', size = 12)
    
    plt.yticks(y_spaced, y_timedelta_str, size = 7)

    # ax.set_ylabel('Time')
    ax.set_title('Average MR SS Times', size = 15)
    ax.legend(loc='upper right')
    plt.grid(axis = 'y', linestyle = '--', linewidth = 0.3)
    if save_fig: plt.savefig('%s/%s' % (image_path, 'avg_mr_ss_times_G11.png'))
    else: plt.show() 
    plt.close()

def graph12(bic, selected_plant, image_path):
    wc_row = []
    ana_month = ''
    # Grab the latest month
    for month in mr_month_dict:
        ana_month = month
    for month in mr_month_dict:
        for plant in mr_month_dict[month]:
            curr_plant = plant.title()
            bic_plant = bic[2].title()
            if curr_plant == bic_plant and month == str(bic[5])[:10]:
                for wc_bic in mr_month_dict[month][plant]:
                    if bic[3] == wc_bic[3]:
                        wc_row.append([curr_plant+ " "+ wc_bic[3], dt2s(wc_bic[40])])
            if curr_plant == selected_plant and month == ana_month: #in the selected plant
                for wc in mr_month_dict[month][plant]:
                    wc_row.append([curr_plant+ " "+ wc[3], dt2s(wc[40])])
    
    presentation_row = []
    for row in wc_row:
        if row[0] not in presentation_row:
            presentation_row.append(row[0])

    #Only use this if I want all months
    tick_list_serial_row = []
    for serial in presentation_row:
        serial_list = []
        for row in wc_row:
            if serial == row[0]:
                # print(row)
                serial_list.append(row)

        combined_serial_list = []
        combined_serial_list.append(serial_list[0][0]) #grab the head
        ot_sum = 0
        for ser in serial_list:
            ot_sum += ser[1]

        combined_serial_list.append(ot_sum)
        tick_list_serial_row.append(combined_serial_list)
    
    #Now we can finally populate the three arrays
    ot_arr = []

    for tick_list in tick_list_serial_row:
        ot_arr.append(tick_list[1])

    tick_array = [ot_arr]

    #Now I can finally plot all these
    for count, tick in enumerate(presentation_row):
        if bic[2] in tick:
            presentation_row[count] = "Best in Class"
    bar12_grapher(presentation_row, tick_array, image_path)
def bar12_grapher(labels, avg_arr, image_path):
    width = 0.35       # the width of the bars: can also be len(x) sequence
    fig, ax = plt.subplots(figsize=(10, 6))
    colors = color_list
    ax.bar(labels, avg_arr[0], width, label='Average Blanket Wash Time', color = colors[0])

    # Y-Tick Factory
    vals = [0 for _ in range(len(avg_arr[0]))]
    for i in range(len(avg_arr[0])): #for WC1
        for metric in avg_arr: 
            vals[i] += metric[i]
    largest_time = max(vals) * 1.1
    y_tick_num = 9
    y_spaced = [largest_time * ((i+1)/y_tick_num) for i in range(y_tick_num)]
    y_timedelta = [str(timedelta(seconds=int(s))) for s in y_spaced]
    if largest_time >= 3600:
        y_timedelta_str = [time_HHMM(s) for s in y_timedelta]
        ax.set_ylabel('Time (h)', size = 12)
    else:
        y_timedelta_str = [time_MMSS(s) for s in y_timedelta]
        ax.set_ylabel('Time (m)', size = 12)
    
    plt.yticks(y_spaced, y_timedelta_str, size = 7)

    # ax.set_ylabel('Time')
    ax.set_title('Average Blanket Wash Time', size = 15)
    ax.legend(loc='upper right')
    plt.grid(axis = 'y', linestyle = '--', linewidth = 0.3)
    if save_fig: plt.savefig('%s/%s' % (image_path, 'avg_blnkt_wash_time_G12.png'))
    else: plt.show() 
    plt.close()

def graph13(bic, selected_plant, image_path):
    wc_row = []
    ana_month = ''
    # Grab the latest month
    for month in mr_month_dict:
        ana_month = month
    for month in mr_month_dict:
        for plant in mr_month_dict[month]:
            curr_plant = plant.title()
            bic_plant = bic[2].title()
            if curr_plant == bic_plant and month == str(bic[5])[:10]:
                for wc_bic in mr_month_dict[month][plant]:
                    if bic[3] == wc_bic[3]:
                        wc_row.append([curr_plant+ " "+ wc_bic[3], wc_bic[13], wc_bic[16]])
            if curr_plant == selected_plant and month == ana_month: #in the selected plant
                for wc in mr_month_dict[month][plant]:
                    wc_row.append([curr_plant+ " "+ wc[3], wc[13], wc[16]])
    
    presentation_row = []
    for row in wc_row:
        if row[0] not in presentation_row:
            presentation_row.append(row[0])

    #Only use this if I want all months
    tick_list_serial_row = []
    for serial in presentation_row:
        serial_list = []
        for row in wc_row:
            if serial == row[0]:
                # print(row)
                serial_list.append(row)

        combined_serial_list = []
        combined_serial_list.append(serial_list[0][0]) #grab the head
        ot_sum, basic_MR, ft_time = 0, 0, 0
        for ser in serial_list:
            ot_sum += ser[1]
            basic_MR += ser[2]
        combined_serial_list.append(ot_sum)
        combined_serial_list.append(basic_MR)
        tick_list_serial_row.append(combined_serial_list)
    
    #Now we can finally populate the three arrays
    ot_arr = []
    basic_MR_arr = []
    for tick_list in tick_list_serial_row:
        ot_arr.append(tick_list[1])
        basic_MR_arr.append(tick_list[2])
    tick_array = [ot_arr, basic_MR_arr]

    #Now I can finally plot all these
    for count, tick in enumerate(presentation_row):
        if bic[2] in tick:
            presentation_row[count] = "Best in Class"
    bar13_grapher(presentation_row, tick_array, image_path)
def bar13_grapher(labels, avg_arr, image_path):
    width = 0.35       # the width of the bars: can also be len(x) sequence
    fig, ax = plt.subplots(figsize=(10, 6))
    colors = color_list
    # print(avg_arr)
    ax.bar(labels, avg_arr[0], width, label='Sum of MR Waste', color = colors[0])
    ax.bar(labels, avg_arr[1], width, bottom=avg_arr[0], 
            label='Sum of Run Waste', color = colors[1])

    ax.set_title('Average Waste', size = 15)
    ax.legend(loc='upper right')
    plt.grid(axis = 'y', linestyle = '--', linewidth = 0.3)
    if save_fig: plt.savefig('%s/%s' % (image_path, 'avg_waste_G13.png'))
    else: plt.show() 
    plt.close()

def graph14(bic, selected_plant, image_path):
    wc_row = []
    ana_month = ''
    # Grab the latest month
    for month in mr_month_dict:
        ana_month = month
    for month in mr_month_dict:
        for plant in mr_month_dict[month]:
            curr_plant = plant.title()
            bic_plant = bic[2].title()
            if curr_plant == bic_plant and month == str(bic[5])[:10]:
                for wc_bic in mr_month_dict[month][plant]:
                    if bic[3] == wc_bic[3]:
                        wc_row.append([curr_plant+ " "+ wc_bic[3], wc_bic[20], 100*wc_bic[20]/wc_bic[18]])
            if curr_plant == selected_plant and month == ana_month: #in the selected plant
                for wc in mr_month_dict[month][plant]:
                    wc_row.append([curr_plant+ " "+ wc[3], wc[20], 100*wc[20]/wc[18]])
    
    presentation_row = []
    for row in wc_row:
        if row[0] not in presentation_row:
            presentation_row.append(row[0])

    #Only use this if I want all months
    tick_list_serial_row = []
    for serial in presentation_row:
        serial_list = []
        for row in wc_row:
            if serial == row[0]:
                # print(row)
                serial_list.append(row)

        combined_serial_list = []
        combined_serial_list.append(serial_list[0][0]) #grab the head
        ot_sum, basic_MR, ft_time = 0, 0, 0
        for ser in serial_list:
            ot_sum += ser[1]
            basic_MR += ser[2]
        combined_serial_list.append(ot_sum)
        combined_serial_list.append(basic_MR)
        tick_list_serial_row.append(combined_serial_list)
    
    #Now we can finally populate the three arrays
    ot_arr = []
    basic_MR_arr = []
    for tick_list in tick_list_serial_row:
        ot_arr.append(tick_list[1])
        basic_MR_arr.append(tick_list[2])
    tick_array = [ot_arr, basic_MR_arr]

    #Now I can finally plot all these
    for count, tick in enumerate(presentation_row):
        if bic[2] in tick:
            presentation_row[count] = "Best in Class"
    bar14_grapher(presentation_row, tick_array, image_path)
def bar14_grapher(labels, avg_arr, image_path):
    width = 0.35       # the width of the bars: can also be len(x) sequence
    fig, ax = plt.subplots(figsize=(10, 6))
    colors = color_list
    
    #Plot the bar charts
    ax.bar(labels, avg_arr[0], width, label='Sum of AVG Net Output', color = colors[0])

    #Set up second axis, plot the markers
    ax2 = ax.twinx()
    plt.plot(labels, avg_arr[1], linestyle='None', marker='D', color='b', label='Sum of AVG Net Output % Of Max')
    
    ax2.set_ylim([0, 100])
    #Labels
    ax.set_ylabel('Net Output', size = 12)
    ax2.set_ylabel('Net Output % Of Max', size = 12) 
    ax.set_title('Average Net Output Speed', size = 15)
    #Legends
    ax.legend(loc='upper right')
    ax2.legend(loc='upper left')
    #grid
    ax.grid(axis = 'y', linestyle = '--', linewidth = 0.3)
    if save_fig: plt.savefig('%s/%s' % (image_path, 'avg_net_output_speed_G14.png'))
    else: plt.show() 
    plt.close()


def subtotal_all(make_ready_excel):
    month_list = []
    df1 = pd.read_excel(make_ready_excel, "GPI location KPI's Mar 2022")
    month_list.append(df1)
    df2 = pd.read_excel(make_ready_excel, "GPI location KPI's Apr 2022")
    month_list.append(df2)
    df3 = pd.read_excel(make_ready_excel, "GPI location KPI's May 2022")
    month_list.append(df3)
    combined = pd.concat(month_list, ignore_index=True)

    
    for index, row in combined.iterrows():
        #grab the month and plant from this row
        curr_month, curr_plant = row["Year Month"].strftime('%Y-%m-%d %X')[:10] , row["Account City"]
        #if the month hasn't been entered yet
        if curr_month not in mr_month_dict:
            mr_month_dict[curr_month] = {}

        #if the plant doesn't exist in month dict yet, add plant
        if curr_plant not in mr_month_dict[curr_month]:
            mr_month_dict[curr_month][curr_plant] = []
        
        #Now, the month and plant MUST exist, so add that row to the correct month-plant bin
        mr_month_dict[curr_month][curr_plant].append(list(row))



    # Find the BIC for the rest...
    bic = []
    smallest_avg_MR = 10**5
    for month in mr_month_dict:
        for plant in mr_month_dict[month]:
            for row in mr_month_dict[month][plant]:
                #27 avg basic MR time, #30 avg Fine Tuning time, #36 avg Other Time
                avg_total_MR_time = dt2s(row[27]) + dt2s(row[30]) + dt2s(row[36])
                if avg_total_MR_time < smallest_avg_MR:
                    bic = row #update the bic
                    smallest_avg_MR = avg_total_MR_time
    # print("The best in class Plant is: ", bic[2])

    analysis_month = ''
    for month in mr_month_dict:
        analysis_month = month


    for plant in plant_list:
        plant_path = os.path.join(long_dir, "plant_report", plant)
        graph_path = os.path.join(plant_path, 'graphs')

        if not os.path.isdir(plant_path):
            os.mkdir(plant_path)
        if not os.path.isdir(graph_path):
            os.mkdir(graph_path)

        if print_graphs:
            graph1(bic, plant.title(), graph_path)
            graph2(bic, plant.title(), graph_path)
            graph3(bic, plant.title(), graph_path)
            graph4(bic, plant.title(), graph_path)
            graph5(bic, plant.title(), graph_path)
            graph6(bic, plant.title(), graph_path)
            graph7(bic, plant.title(), graph_path)
            graph8(bic, plant.title(), graph_path)
            graph9(bic, plant.title(), graph_path)
            graph10(bic, plant.title(), graph_path)
            graph11(bic, plant.title(), graph_path)
            graph12(bic, plant.title(), graph_path)
            graph13(bic, plant.title(), graph_path)
            graph14(bic, plant.title(), graph_path)
        

        create_ppt(plant, plant_path, analysis_month)







def create_ppt(plant, plant_dir, curr_month):
    month = ('%s_%s') % (month_label[int(curr_month[5:7])-3], curr_month[:4])
    #Create image list so I can access the pictures easily
    graph_path = os.path.join(plant_dir, 'graphs')
    image_list = []
    image_list.append('index_0_G0.png')
    for i in range(14):
        for file in os.listdir(graph_path):
            if i+1 == int(file[file.find('_G')+2:file.find('.png')]):
                image_list.append(file)

    prs = Presentation('gpi_template.pptx')



    ############TITLE SLIDE #################
    #Choode a slide layout
    lyt=prs.slide_layouts[0]
    slide=prs.slides.add_slide(lyt) # adding a slide

    #Just looking for the place holder names (optional)
    # for shape in slide.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))

    #title for this slide
    title=slide.shapes.title # assigning a title
    title_name= ("Heidelberg Make Ready Time Analysis For %s - %s") % (plant, month.replace('_', ' ')) #title
    text_ph1 = slide.placeholders[0]
    text_ph1.text_frame.text = title_name
    font = text_ph1.text_frame.paragraphs[0].runs[0].font
    font.name = 'Arial'
    font.size = Pt(29)
    font.bold = True
    font.color.rgb = RGBColor(0x00, 0x00, 0x00) #change the title color
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER #Center the title

    #Delete the subtitle place holder
    subtitle=slide.placeholders[1]
    sp = subtitle.element
    sp.getparent().remove(sp)



    ############ Terminology SLIDE #################
    term_slide = prs.slide_layouts[3]
    slide=prs.slides.add_slide(term_slide) # adding a slide
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Terminology'

    tf = body_shape.text_frame
    tf.text = 'Each of these metrics (time spans) also has an associated stand still time (SS).  Stand still time is the time where the press is receiving no commands from the console or running any programs (processes)\n'

    p = tf.add_paragraph()
    p.text = 'Other Time'
    p.level = 1
    p.runs[0].font.bold = True

    p = tf.add_paragraph()
    p.text = 'Time from last good sheet to start of plate changing'
    p.level = 2


    p = tf.add_paragraph()
    p.text = 'Basic Make Ready Time'
    p.level = 1
    p.runs[0].font.bold = True
    p = tf.add_paragraph()
    p.text = 'Time from start of plate changing to first printed sheet'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Fine Tune Time'
    p.level = 1
    p.runs[0].font.bold = True
    p = tf.add_paragraph()
    p.text = 'Time from first printed sheet to first good sheet (registration and color Ok)'
    p.level = 2

    p = tf.add_paragraph()
    p.text = 'Total Make Ready Time'
    p.level = 1
    p.runs[0].font.bold = True
    p = tf.add_paragraph()
    p.text = 'Other Time + Basic Make Ready Time + Fine Tune Time'
    p.level = 2
    p = tf.add_paragraph()
    p.text = 'This equates to time from last good sheet to first good sheet on next job'
    p.level = 2


    p = tf.add_paragraph()
    p.text = 'Best In Class'
    p.level = 1
    p.runs[0].font.bold = True
    p = tf.add_paragraph()
    p.text = 'Baseline set for comparison'
    p.level = 2
    p = tf.add_paragraph()
    p.text = 'Determined by the lowest “Avg Total Make Ready Time”'
    p.level = 2
    









    ############ Total MR Time [3] #################
    slide_G1 = prs.slide_layouts[6]
    slide_G1s = prs.slides.add_slide(slide_G1) # adding a slide
    shapes_G1 = slide_G1s.shapes
    title_shape_mr = shapes_G1.title
    title_shape_mr.text = 'Total Make Ready Time'
    # for shape in slide_G1s.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    img_path = os.path.join(graph_path, image_list[1])

    pic_holder = slide_G1s.placeholders[13]
    # print(pic_holder.shape_type)
    # print(pic_holder.name)
    placeholder_picture = pic_holder.insert_picture(img_path)
    body_shape = shapes_G1.placeholders[14]
    tf = body_shape.text_frame
    tf.text = 'Total Make Ready Time = Other Time + Basic Make Ready Time + Fine Tune Time Last good sheet to first good sheet'





    ############ Total MR Time [4] #################
    slide_G2 = prs.slide_layouts[6]
    slide_G2s = prs.slides.add_slide(slide_G2) # adding a slide
    shapes_G2 = slide_G2s.shapes
    title_shape_mr = shapes_G2.title
    title_shape_mr.text = 'Make Ready Time as Lost Sheets'
    # for shape in slide_G2s.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    img_path = os.path.join(graph_path, image_list[2])

    pic_holder = slide_G2s.placeholders[13]
    placeholder_picture = pic_holder.insert_picture(img_path)
    body_shape = shapes_G2.placeholders[14]
    tf = body_shape.text_frame
    tf.text = 'Total Make Ready Time = Other Time + Basic Make Ready Time + Fine Tune Time Last good sheet to first good sheet'



    ############ Total MR Time [5] #################
    slide_G3 = prs.slide_layouts[6]
    slide_G3s = prs.slides.add_slide(slide_G3) # adding a slide
    shapes_G3 = slide_G3s.shapes
    title_shape_mr = shapes_G3.title
    title_shape_mr.text = 'Total Make Ready Time'
    # for shape in slide_G3s.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    img_path = os.path.join(graph_path, image_list[3])

    pic_holder = slide_G3s.placeholders[13]
    placeholder_picture = pic_holder.insert_picture(img_path)
    body_shape = shapes_G3.placeholders[14]
    tf = body_shape.text_frame
    tf.text = 'Total Make Ready Time = Other Time + Basic Make Ready Time + Fine Tune Time Last good sheet to first good sheet'


    ############ Total MR Time [6] #################
    slide_G4 = prs.slide_layouts[6]
    slide_G4s = prs.slides.add_slide(slide_G4) # adding a slide
    shapes_G4 = slide_G4s.shapes
    title_shape_mr = shapes_G4.title
    title_shape_mr.text = 'Make Ready Time as Lost Sheets'
    # for shape in slide_G4s.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    img_path = os.path.join(graph_path, image_list[4])

    pic_holder = slide_G4s.placeholders[13]
    placeholder_picture = pic_holder.insert_picture(img_path)
    body_shape = shapes_G4.placeholders[14]
    tf = body_shape.text_frame
    tf.text = 'Total Make Ready Time = Other Time + Basic Make Ready Time + Fine Tune Time Last good sheet to first good sheet'


    ############ Total MR Time [7] #################
    slide_G5 = prs.slide_layouts[6]
    slide_G5s = prs.slides.add_slide(slide_G5) # adding a slide
    shapes_G5 = slide_G5s.shapes
    title_shape_mr = shapes_G5.title
    title_shape_mr.text = 'Other Make Ready Time'
    # for shape in slide_G5s.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    img_path = os.path.join(graph_path, image_list[5])

    pic_holder = slide_G5s.placeholders[13]
    placeholder_picture = pic_holder.insert_picture(img_path)
    body_shape = shapes_G5.placeholders[14]
    tf = body_shape.text_frame
    tf.text = 'Other Time = Time from last good sheet to start of plate changing'

    ############ Total MR Time [8] #################
    slide_G6 = prs.slide_layouts[6]
    slide_G6s = prs.slides.add_slide(slide_G6) # adding a slide
    shapes_G6 = slide_G6s.shapes
    title_shape_mr = shapes_G6.title
    title_shape_mr.text = 'Other Make Ready Time'
    # for shape in slide_G6s.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    img_path = os.path.join(graph_path, image_list[6])

    pic_holder = slide_G6s.placeholders[13]
    placeholder_picture = pic_holder.insert_picture(img_path)
    body_shape = shapes_G6.placeholders[14]
    tf = body_shape.text_frame
    tf.text = 'Other Time = Time from last good sheet to start of plate changing'

    ############ Total MR Time [9] #################
    slide_G7 = prs.slide_layouts[6]
    slide_G7s = prs.slides.add_slide(slide_G7) # adding a slide
    shapes_G7 = slide_G7s.shapes
    title_shape_mr = shapes_G7.title
    title_shape_mr.text = 'Basic Make Ready Time'
    # for shape in slide_G7s.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    img_path = os.path.join(graph_path, image_list[7])

    pic_holder = slide_G7s.placeholders[13]
    placeholder_picture = pic_holder.insert_picture(img_path)
    body_shape = shapes_G7.placeholders[14]
    tf = body_shape.text_frame
    tf.text = 'Basic Make Ready Time =Time from start of plate changing to first printed sheet'

    ############ Total MR Time [10] #################
    slide_G8 = prs.slide_layouts[6]
    slide_G8s = prs.slides.add_slide(slide_G8) # adding a slide
    shapes_G8 = slide_G8s.shapes
    title_shape_mr = shapes_G8.title
    title_shape_mr.text = 'Basic Make Ready Time as Lost Sheets'
    # for shape in slide_G8s.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    img_path = os.path.join(graph_path, image_list[8])

    pic_holder = slide_G8s.placeholders[13]
    placeholder_picture = pic_holder.insert_picture(img_path)
    body_shape = shapes_G8.placeholders[14]
    tf = body_shape.text_frame
    tf.text = 'Basic Make Ready Time =Time from start of plate changing to first printed sheet'

    ############ Total MR Time [11] #################
    slide_G9 = prs.slide_layouts[6]
    slide_G9s = prs.slides.add_slide(slide_G9) # adding a slide
    shapes_G9 = slide_G9s.shapes
    title_shape_mr = shapes_G9.title
    title_shape_mr.text = 'Fine Tune Time'
    # for shape in slide_G9s.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    img_path = os.path.join(graph_path, image_list[9])

    pic_holder = slide_G9s.placeholders[13]
    placeholder_picture = pic_holder.insert_picture(img_path)
    body_shape = shapes_G9.placeholders[14]
    tf = body_shape.text_frame
    tf.text = 'Fine Tune Time = Time from first printed sheet to first good sheet (registration and color Ok)'

    ############ Total MR Time [12] #################
    slide_G10 = prs.slide_layouts[6]
    slide_G10s = prs.slides.add_slide(slide_G10) # adding a slide
    shapes_G10 = slide_G10s.shapes
    title_shape_mr = shapes_G10.title
    title_shape_mr.text = 'Fine Tune Time as Lost Sheets'
    # for shape in slide_G10s.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    img_path = os.path.join(graph_path, image_list[10])

    pic_holder = slide_G10s.placeholders[13]
    placeholder_picture = pic_holder.insert_picture(img_path)
    body_shape = shapes_G10.placeholders[14]
    tf = body_shape.text_frame
    tf.text = 'Fine Tune Time = Time from first printed sheet to first good sheet (registration and color Ok)'

    ############ Total MR Time [13] #################
    slide_G11 = prs.slide_layouts[6]
    slide_G11s = prs.slides.add_slide(slide_G11) # adding a slide
    shapes_G11 = slide_G11s.shapes
    title_shape_mr = shapes_G11.title
    title_shape_mr.text = 'Stand Still Times'
    # for shape in slide_G11s.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    img_path = os.path.join(graph_path, image_list[11])

    pic_holder = slide_G11s.placeholders[13]
    placeholder_picture = pic_holder.insert_picture(img_path)
    body_shape = shapes_G11.placeholders[14]
    tf = body_shape.text_frame
    tf.text = 'Stand still time is the time where the press is receiving no commands from the console or running any programs (processes)'





    ############ Additional Metrics SLIDE [14] #################
    am_slide = prs.slide_layouts[3]
    slide=prs.slides.add_slide(am_slide) # adding a slide
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Additional Metrics'

    tf = body_shape.text_frame
    
    p = tf.add_paragraph()
    p.text = '\n\n\nAdditional Metrics Affecting Overall Throughput'
    p.runs[0].font.bold = True
    p.runs[0].font.name = 'Calibri'
    p.runs[0].font.size = Pt(28)
    p.alignment = PP_ALIGN.CENTER #Center the title



    ############ Total MR Time [15] #################
    slide_G12 = prs.slide_layouts[7]
    slide_G12s = prs.slides.add_slide(slide_G12) # adding a slide
    shapes_G12 = slide_G12s.shapes
    title_shape_mr = shapes_G12.title
    title_shape_mr.text = 'Blanket Washing Time'
    # for shape in slide_G12s.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    img_path = os.path.join(graph_path, image_list[12])

    pic_holder = slide_G12s.placeholders[13]
    pic_holder.insert_picture(img_path)


    ############ Total MR Time [17] #################
    slide_G13 = prs.slide_layouts[7]
    slide_G13s = prs.slides.add_slide(slide_G13) # adding a slide
    shapes_G13 = slide_G13s.shapes
    title_shape_mr = shapes_G13.title
    title_shape_mr.text = 'Blanket Washing Time'
    # for shape in slide_G13s.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    img_path = os.path.join(graph_path, image_list[13])

    pic_holder = slide_G13s.placeholders[13]
    pic_holder.insert_picture(img_path)



    ############ Total MR Time [18] #################
    slide_G14 = prs.slide_layouts[7]
    slide_G14s = prs.slides.add_slide(slide_G14) # adding a slide
    shapes_G14 = slide_G14s.shapes
    title_shape_mr = shapes_G14.title
    title_shape_mr.text = 'Blanket Washing Time'
    # for shape in slide_G14s.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    img_path = os.path.join(graph_path, image_list[14])

    pic_holder = slide_G14s.placeholders[13]
    pic_holder.insert_picture(img_path)













    ############ Last Slide Questions [18] #################
    am_slide = prs.slide_layouts[3]
    slide=prs.slides.add_slide(am_slide) # adding a slide
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = '7204 Overview'

    tf = body_shape.text_frame
    
    p = tf.add_paragraph()
    p.text = '\n\n\n\t\t\t\t\tQuestions?' #hard coded formatting lol
    p.runs[0].font.bold = True
    p.runs[0].font.name = 'Calibri'
    p.runs[0].font.size = Pt(54)



    ############ SAVING THE FILE #################
    filename = ('HB_NA_Fleet_Make_Ready_Data_%s_%s.pptx') % (plant, month)
    report_path = os.path.join(plant_dir, 'report')
    if not os.path.isdir(report_path): os.mkdir(report_path)
    file_dest = os.path.join(report_path, filename)
    prs.save(file_dest) # saving file

    #Save a copy in the base file month as well
    base_report_path = os.path.join(long_dir, 'report_folder', curr_month)
    if not os.path.isdir(base_report_path): os.mkdir(base_report_path)
    shutil.copy(file_dest, base_report_path)






def main():
    #load excel database into a pd format
    make_ready_excel = pd.ExcelFile(current)
    # throw the pd table into the function to figure out

    subtotal_all(make_ready_excel)




if __name__ == '__main__':
    main()



'''


    if row['Die Description'] in die_desc_dict_may: # if this description is in the current month
        total_time = row['Setup']+row['Run']+row['Down'] 
        if total_time > 0 and row['Confirmed Qty'] > 0: # Make sure time actually exists before I do anything
            

            # if time exists, input the WC, Q, Rate, and Index
            # row_info = [row['Work Center'], total_time,row['Confirmed Qty'], row['Confirmed Qty']/total_time, index]
            # if a new Die Desc, then make an array to hold the info above
            if row['Die Description'] not in die_dict: #if it's a new Die, 
                die_dict[row['Die Description']] = []
                die_dict[row['Die Description']].append(list(row))
            # if already a Die desc, then append the current info
            else:
                die_dict[row['Die Description']].append(list(row))


# This is for creating new columns in excel, right now I'm not doing that yet
# X.loc[:, "time_delta"] = time_diff_list
# X.loc[:, "Sheet Diff"] = net_diff
# print("Total number of descriptions:", len(desc_list))


repeat_checker = 0
repeat_wc_key_buffer = []
for key, die_order in die_dict.items():
    # print(key, die_order)

    for work_order in die_order:
        print(work_order[0], round(work_order[3],2))

        if repeat_checker:
            if work_order[0] in repeat_wc_key_buffer: print("DOUBLE FOUND")
            else: repeat_wc_key_buffer.append(work_order[0])
    
    # clear out the key buffer
    repeat_wc_key_buffer = []

    # My goal is to create a 2021 values of die description, with a "max" value for each 


# Now I will use this die dictionary to comb through the May graph
for index, row in Y.iterrows(): # check every row in 2021 year
    if row['Die Description'] in die_desc_dict_may:
        # time_string = row['Date'].strftime('%Y-%m-%d %X')
        # print(time_string[:10])
        print(list(row))










#### BAR 1 EXAMPLE ############
def bar1_example(X, Y, Z, A):
    X = ['Best Asset','WC1','WC2', 'WC3'] #Name of Plant + WC
    Y = [10,25,32, 10] #March
    Z = [24,34,10, 12] #April
    A = [57, 43, 23, 15] #May
    bar1_grapher(X, [Y,Z,A])



'''
